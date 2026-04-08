"""
PPTX Parser — Extracts text, tables, images, and metadata from PowerPoint files.

Returns structured `PresentationContent` objects ready for the AI extraction layer.
"""

from __future__ import annotations

import base64
import io
import logging
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches

logger = logging.getLogger(__name__)


# ── Data Classes ────────────────────────────────────────────────────────────


@dataclass
class TableData:
    """Represents a single table extracted from a slide."""

    headers: list[str]
    rows: list[dict[str, str]]

    def to_markdown(self) -> str:
        """Render the table as a Markdown string for prompt inclusion."""
        if not self.headers:
            return ""
        header_line = "| " + " | ".join(self.headers) + " |"
        separator = "| " + " | ".join(["---"] * len(self.headers)) + " |"
        body_lines = []
        for row in self.rows:
            cells = [row.get(h, "") for h in self.headers]
            body_lines.append("| " + " | ".join(cells) + " |")
        return "\n".join([header_line, separator, *body_lines])


@dataclass
class ImageData:
    """Represents an extracted image from a slide."""

    name: str
    content_type: str
    base64_data: str
    width_inches: float | None = None
    height_inches: float | None = None


@dataclass
class SlideContent:
    """Structured content extracted from a single slide."""

    slide_number: int
    title: str | None = None
    layout_name: str | None = None
    text_blocks: list[str] = field(default_factory=list)
    tables: list[TableData] = field(default_factory=list)
    images: list[ImageData] = field(default_factory=list)
    speaker_notes: str | None = None

    @property
    def has_content(self) -> bool:
        """Check if the slide has any meaningful content."""
        return bool(
            self.text_blocks
            or self.tables
            or self.speaker_notes
        )

    def to_text(self, include_images: bool = False) -> str:
        """
        Render all slide content as a formatted text block
        suitable for prompt injection.
        """
        parts: list[str] = []
        parts.append(f"=== Slide {self.slide_number} ===")

        if self.title:
            parts.append(f"Title: {self.title}")
        if self.layout_name:
            parts.append(f"Layout: {self.layout_name}")

        if self.text_blocks:
            parts.append("\n--- Text Content ---")
            for block in self.text_blocks:
                stripped = block.strip()
                if stripped:
                    parts.append(stripped)

        if self.tables:
            parts.append("\n--- Tables ---")
            for i, table in enumerate(self.tables, 1):
                parts.append(f"Table {i}:")
                parts.append(table.to_markdown())

        if self.speaker_notes:
            parts.append(f"\n--- Speaker Notes ---\n{self.speaker_notes}")

        if include_images and self.images:
            parts.append(f"\n[{len(self.images)} image(s) on this slide]")

        return "\n".join(parts)


@dataclass
class PresentationContent:
    """Aggregated content from an entire presentation."""

    file_path: str
    file_name: str
    total_slides: int
    slides: list[SlideContent] = field(default_factory=list)
    metadata: dict[str, Any] = field(default_factory=dict)

    @property
    def content_slides(self) -> list[SlideContent]:
        """Return only slides that contain meaningful content."""
        return [s for s in self.slides if s.has_content]

    def to_full_text(self, include_images: bool = False) -> str:
        """Render the entire presentation as a single text document."""
        header = (
            f"Presentation: {self.file_name}\n"
            f"Total Slides: {self.total_slides}\n"
            f"Content Slides: {len(self.content_slides)}\n"
            f"{'=' * 60}\n"
        )
        slide_texts = [
            s.to_text(include_images=include_images) for s in self.content_slides
        ]
        return header + "\n\n".join(slide_texts)


# ── Parser ──────────────────────────────────────────────────────────────────


class PPTXParser:
    """
    Parses a .pptx file and extracts structured content from each slide.

    Usage:
        parser = PPTXParser()
        content = parser.parse("path/to/deck.pptx")
        print(content.to_full_text())
    """

    def parse(self, file_path: str | Path) -> PresentationContent:
        """
        Parse a .pptx file and return structured content.

        Args:
            file_path: Path to the PowerPoint file.

        Returns:
            PresentationContent with all extracted data.

        Raises:
            FileNotFoundError: If the file does not exist.
            ValueError: If the file is not a .pptx file.
        """
        path = Path(file_path)

        if not path.exists():
            raise FileNotFoundError(f"File not found: {path}")
        if path.suffix.lower() != ".pptx":
            raise ValueError(f"Expected .pptx file, got: {path.suffix}")

        logger.info("Parsing presentation: %s", path.name)

        prs = Presentation(str(path))
        slides: list[SlideContent] = []

        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_content = self._extract_slide(slide, slide_num)
            slides.append(slide_content)

        content = PresentationContent(
            file_path=str(path.resolve()),
            file_name=path.name,
            total_slides=len(prs.slides),
            slides=slides,
            metadata=self._extract_metadata(prs),
        )

        logger.info(
            "Parsed %d slides (%d with content) from %s",
            content.total_slides,
            len(content.content_slides),
            path.name,
        )
        return content

    # ── Private Helpers ─────────────────────────────────────────────────

    def _extract_slide(self, slide, slide_num: int) -> SlideContent:
        """Extract all content from a single slide."""
        content = SlideContent(slide_number=slide_num)

        # Layout name
        try:
            content.layout_name = slide.slide_layout.name
        except Exception:
            pass

        # Title
        if slide.shapes.title and slide.shapes.title.text:
            content.title = slide.shapes.title.text.strip()

        # Iterate all shapes
        for shape in slide.shapes:
            # Skip the title shape (already captured)
            if shape == slide.shapes.title:
                continue

            # Tables
            if shape.has_table:
                table_data = self._extract_table(shape.table)
                if table_data:
                    content.tables.append(table_data)

            # Images
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image_data = self._extract_image(shape)
                if image_data:
                    content.images.append(image_data)

            # Text frames
            elif shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    content.text_blocks.append(text)

            # Group shapes — recurse
            elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                self._extract_group(shape, content)

        # Speaker notes
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                content.speaker_notes = notes

        return content

    def _extract_table(self, table) -> TableData | None:
        """Extract table data as headers + rows."""
        try:
            rows = list(table.rows)
            if not rows:
                return None

            # First row as headers
            headers = [cell.text.strip() for cell in rows[0].cells]

            # Remaining rows as data
            data_rows = []
            for row in rows[1:]:
                cells = [cell.text.strip() for cell in row.cells]
                row_dict = dict(zip(headers, cells))
                data_rows.append(row_dict)

            return TableData(headers=headers, rows=data_rows)

        except Exception as e:
            logger.warning("Failed to extract table: %s", e)
            return None

    def _extract_image(self, shape) -> ImageData | None:
        """Extract image data as base64."""
        try:
            image = shape.image
            blob = image.blob
            content_type = image.content_type

            b64 = base64.b64encode(blob).decode("utf-8")

            width = None
            height = None
            if shape.width:
                width = round(shape.width / Inches(1), 2)
            if shape.height:
                height = round(shape.height / Inches(1), 2)

            return ImageData(
                name=shape.name or "unnamed",
                content_type=content_type,
                base64_data=b64,
                width_inches=width,
                height_inches=height,
            )

        except Exception as e:
            logger.warning("Failed to extract image '%s': %s", shape.name, e)
            return None

    def _extract_group(self, group_shape, content: SlideContent) -> None:
        """Recursively extract content from grouped shapes."""
        try:
            for shape in group_shape.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if text:
                        content.text_blocks.append(text)
                elif shape.has_table:
                    table_data = self._extract_table(shape.table)
                    if table_data:
                        content.tables.append(table_data)
                elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image_data = self._extract_image(shape)
                    if image_data:
                        content.images.append(image_data)
        except Exception as e:
            logger.warning("Failed to extract group shape: %s", e)

    def _extract_metadata(self, prs: Presentation) -> dict[str, Any]:
        """Extract presentation-level metadata."""
        meta: dict[str, Any] = {}
        try:
            core = prs.core_properties
            if core.title:
                meta["title"] = core.title
            if core.author:
                meta["author"] = core.author
            if core.subject:
                meta["subject"] = core.subject
            if core.created:
                meta["created"] = core.created.isoformat()
            if core.modified:
                meta["modified"] = core.modified.isoformat()
            if core.last_modified_by:
                meta["last_modified_by"] = core.last_modified_by
        except Exception as e:
            logger.warning("Failed to extract metadata: %s", e)
        return meta
